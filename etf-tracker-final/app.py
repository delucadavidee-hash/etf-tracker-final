# ═══════════════════════════════════════════
# ETF Tracker - Flask Backend with SQLite
# Plug & Play - zero config, single file DB
# ═══════════════════════════════════════════

import os, sqlite3, json, uuid, hashlib, secrets, smtplib
from datetime import datetime, timedelta
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders
from functools import wraps
from flask import Flask, request, jsonify, render_template, send_file, g
from flask_cors import CORS

# ─── Config ───
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DB_PATH = os.path.join(BASE_DIR, 'data_store.json')
UPLOAD_FOLDER = os.path.join(BASE_DIR, 'reports')
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

SECRET_KEY = os.environ.get('SECRET_KEY', secrets.token_hex(32))
SMTP_HOST = os.environ.get('SMTP_HOST', 'smtp.gmail.com')
SMTP_PORT = int(os.environ.get('SMTP_PORT', 587))
SMTP_USER = os.environ.get('SMTP_USER', '')
SMTP_PASS = os.environ.get('SMTP_PASS', '')
FROM_EMAIL = os.environ.get('FROM_EMAIL', SMTP_USER)

# ─── App ───
app = Flask(__name__, template_folder='templates', static_folder='static')
app.config['SECRET_KEY'] = SECRET_KEY
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024
CORS(app)

# ─── Database ───
def get_db():
    db = getattr(g, '_database', None)
    if db is None:
        db = g._database = sqlite3.connect(DB_PATH)
        db.row_factory = sqlite3.Row
    return db

def init_db():
    db = get_db()
    cursor = db.cursor()
    
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            email TEXT UNIQUE NOT NULL,
            password_hash TEXT NOT NULL,
            name TEXT DEFAULT 'Investitore',
            initials TEXT DEFAULT 'IN',
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            last_login TIMESTAMP,
            onboarding_done INTEGER DEFAULT 0,
            onboarding_answers TEXT
        )
    ''')
    
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS portfolios (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            name TEXT DEFAULT 'Portafoglio Principale',
            strategy TEXT DEFAULT 'custom',
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (user_id) REFERENCES users(id) ON DELETE CASCADE
        )
    ''')
    
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS holdings (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            portfolio_id INTEGER NOT NULL,
            isin TEXT NOT NULL,
            ticker TEXT NOT NULL,
            name TEXT,
            shares REAL NOT NULL DEFAULT 0,
            avg_price REAL NOT NULL DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (portfolio_id) REFERENCES portfolios(id) ON DELETE CASCADE
        )
    ''')
    
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS alerts (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            etf_isin TEXT NOT NULL,
            etf_ticker TEXT NOT NULL,
            alert_type TEXT NOT NULL,
            threshold REAL NOT NULL,
            current_price REAL,
            active INTEGER DEFAULT 1,
            channels TEXT DEFAULT '["email","push"]',
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (user_id) REFERENCES users(id) ON DELETE CASCADE
        )
    ''')
    
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS settings (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            key TEXT NOT NULL,
            value TEXT,
            UNIQUE(user_id, key),
            FOREIGN KEY (user_id) REFERENCES users(id) ON DELETE CASCADE
        )
    ''')
    
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS reports (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            report_type TEXT NOT NULL,
            format TEXT NOT NULL,
            title TEXT,
            content_json TEXT,
            file_path TEXT,
            sent_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            status TEXT DEFAULT 'generated',
            FOREIGN KEY (user_id) REFERENCES users(id) ON DELETE CASCADE
        )
    ''')
    
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS community_posts (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            user_name TEXT NOT NULL,
            avatar TEXT,
            content TEXT NOT NULL,
            likes INTEGER DEFAULT 0,
            comments INTEGER DEFAULT 0,
            portfolio_snapshot TEXT,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    db.commit()
    seed_data(cursor, db)

def seed_data(cursor, db):
    # Seed community posts if empty
    cursor.execute('SELECT COUNT(*) FROM community_posts')
    if cursor.fetchone()[0] == 0:
        posts = [
            ('ETF Enthusiast', 'ET', 'Ho appena ribilanciato il mio portafoglio All-Weather dopo il calo dei bond. Ora 40% azionario, 30% obbligazionario, 15% oro, 15% commodities. CAGR atteso 7.2%.', 12, 3, '{"cagr":7.2,"assets":4}'),
            ('Boglehead Italia', 'BI', 'Dopo 3 anni di PAC su VWCE, ho finalmente superato la soglia dei 50k. Disciplina e pazienza battono il timing del mercato.', 8, 1, None),
            ('Gold Bug', 'GB', 'Aggiunto 5% oro fisico via Xetra-Gold al mio portafoglio. Non per rendimento, ma per drawdown protection.', 5, 0, None),
        ]
        for name, av, content, likes, comments, snap in posts:
            cursor.execute('''
                INSERT INTO community_posts (user_name, avatar, content, likes, comments, portfolio_snapshot)
                VALUES (?, ?, ?, ?, ?, ?)
            ''', (name, av, content, likes, comments, snap))
        db.commit()

@app.teardown_appcontext
def close_db(exception):
    db = getattr(g, '_database', None)
    if db is not None:
        db.close()

# ─── Auth Helpers ───
def hash_password(password):
    return hashlib.pbkdf2_hmac('sha256', password.encode(), SECRET_KEY.encode(), 100000).hex()

def generate_token(user_id):
    payload = {'user_id': user_id, 'exp': datetime.utcnow() + timedelta(days=7), 'jti': str(uuid.uuid4())}
    # Simple token: base64(json) format for stateless auth
    import base64
    token_data = base64.b64encode(json.dumps(payload, default=str).encode()).decode()
    sig = hashlib.sha256((token_data + SECRET_KEY).encode()).hexdigest()[:32]
    return f"{token_data}.{sig}"

def verify_token(token):
    try:
        parts = token.split('.')
        if len(parts) != 2: return None
        token_data, sig = parts
        expected = hashlib.sha256((token_data + SECRET_KEY).encode()).hexdigest()[:32]
        if sig != expected: return None
        import base64
        payload = json.loads(base64.b64decode(token_data.encode()))
        exp = datetime.fromisoformat(payload['exp'].replace('Z', '+00:00').replace('+00:00', ''))
        if datetime.utcnow() > exp: return None
        return payload['user_id']
    except Exception:
        return None

def token_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        token = request.headers.get('Authorization', '').replace('Bearer ', '')
        if not token:
            return jsonify({'error': 'Token mancante'}), 401
        user_id = verify_token(token)
        if not user_id:
            return jsonify({'error': 'Token non valido o scaduto'}), 401
        g.current_user_id = user_id
        return f(*args, **kwargs)
    return decorated

def get_current_user():
    db = get_db()
    cursor = db.cursor()
    cursor.execute('SELECT * FROM users WHERE id = ?', (g.current_user_id,))
    row = cursor.fetchone()
    if row:
        return dict(row)
    return None

# ─── Email Service ───
def send_email(to_email, subject, body_html, attachments=None):
    if not all([SMTP_USER, SMTP_PASS]):
        print(f"[EMAIL] Configurazione SMTP mancante. Email a {to_email} non inviata.")
        return False
    try:
        msg = MIMEMultipart('alternative')
        msg['Subject'] = subject
        msg['From'] = FROM_EMAIL
        msg['To'] = to_email
        msg.attach(MIMEText(body_html, 'html'))
        
        if attachments:
            for filepath in attachments:
                if os.path.exists(filepath):
                    with open(filepath, 'rb') as f:
                        part = MIMEBase('application', 'octet-stream')
                        part.set_payload(f.read())
                    encoders.encode_base64(part)
                    part.add_header('Content-Disposition', f'attachment; filename="{os.path.basename(filepath)}"')
                    msg.attach(part)
        
        with smtplib.SMTP(SMTP_HOST, SMTP_PORT) as server:
            server.starttls()
            server.login(SMTP_USER, SMTP_PASS)
            server.sendmail(FROM_EMAIL, [to_email], msg.as_string())
        print(f"[EMAIL] Inviata a {to_email}: {subject}")
        return True
    except Exception as e:
        print(f"[EMAIL] Errore: {e}")
        return False

# ─── Report Generators ───
def generate_pdf_report(user_id, portfolio_data, report_type='portfolio_summary'):
    """Generate PDF from HTML template using WeasyPrint if available, else reportlab"""
    try:
        from weasyprint import HTML, CSS
        html_content = build_report_html(user_id, portfolio_data, report_type)
        pdf_path = os.path.join(UPLOAD_FOLDER, f'report_{user_id}_{int(datetime.now().timestamp())}.pdf')
        HTML(string=html_content).write_pdf(pdf_path)
        return pdf_path
    except ImportError:
        # Fallback: use reportlab
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
            from reportlab.lib import colors
            pdf_path = os.path.join(UPLOAD_FOLDER, f'report_{user_id}_{int(datetime.now().timestamp())}.pdf')
            doc = SimpleDocTemplate(pdf_path, pagesize=A4)
            styles = getSampleStyleSheet()
            story = []
            
            user = get_current_user() or {'name': 'Investitore', 'email': ''}
            story.append(Paragraph(f"<b>Riepilogo Portafoglio - {user['name']}</b>", styles['Heading1']))
            story.append(Spacer(1, 20))
            story.append(Paragraph(f"Generato il: {datetime.now().strftime('%d/%m/%Y %H:%M')}", styles['Normal']))
            story.append(Spacer(1, 20))
            
            if portfolio_data:
                total_value = sum(h['shares'] * h.get('current_price', h['avg_price']) for h in portfolio_data)
                invested = sum(h['shares'] * h['avg_price'] for h in portfolio_data)
                pl = total_value - invested
                
                story.append(Paragraph(f"<b>Valore Totale:</b> EUR {total_value:,.2f}", styles['Heading2']))
                story.append(Paragraph(f"<b>Investito:</b> EUR {invested:,.2f}", styles['Normal']))
                story.append(Paragraph(f"<b>P&L:</b> EUR {pl:,.2f} ({pl/invested*100 if invested else 0:.1f}%)", styles['Normal']))
                story.append(Spacer(1, 20))
                
                table_data = [['ETF', 'Quote', 'Prezzo Medio', 'Prezzo Att', 'Valore', 'P&L']]
                for h in portfolio_data:
                    cp = h.get('current_price', h['avg_price'])
                    val = h['shares'] * cp
                    plh = h['shares'] * (cp - h['avg_price'])
                    table_data.append([
                        h['ticker'],
                        f"{h['shares']:.2f}",
                        f"{h['avg_price']:.2f}",
                        f"{cp:.2f}",
                        f"{val:,.2f}",
                        f"{plh:,.2f}"
                    ])
                
                t = Table(table_data)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#0A2540')),
                    ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
                    ('ALIGN', (0,0), (-1,-1), 'CENTER'),
                    ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0,0), (-1,0), 11),
                    ('BOTTOMPADDING', (0,0), (-1,0), 12),
                    ('BACKGROUND', (0,1), (-1,-1), colors.HexColor('#F7F9FC')),
                    ('GRID', (0,0), (-1,-1), 1, colors.HexColor('#E8ECF0')),
                ]))
                story.append(t)
            
            doc.build(story)
            return pdf_path
        except ImportError:
            # Ultimate fallback: create a simple text file
            txt_path = os.path.join(UPLOAD_FOLDER, f'report_{user_id}_{int(datetime.now().timestamp())}.txt')
            with open(txt_path, 'w') as f:
                f.write(f"Riepilogo Portafoglio\n")
                f.write(f"Generato: {datetime.now().strftime('%d/%m/%Y %H:%M')}\n\n")
                if portfolio_data:
                    for h in portfolio_data:
                        f.write(f"{h['ticker']}: {h['shares']} quote @ {h['avg_price']:.2f}\n")
            return txt_path

def generate_pptx_report(user_id, portfolio_data, report_type='portfolio_summary'):
    """Generate PPTX using python-pptx"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(6, 21, 40)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "ETF Tracker - Riepilogo Portafoglio"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        user = get_current_user() or {'name': 'Investitore'}
        subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12), Inches(0.5))
        stf = subtitle.text_frame
        sp = stf.paragraphs[0]
        sp.text = f"Generato per {user['name']} - {datetime.now().strftime('%d/%m/%Y')}"
        sp.font.size = Pt(18)
        sp.font.color.rgb = RGBColor(134, 239, 172)
        sp.alignment = PP_ALIGN.CENTER
        
        # KPI slide
        if portfolio_data:
            total_value = sum(h['shares'] * h.get('current_price', h['avg_price']) for h in portfolio_data)
            invested = sum(h['shares'] * h['avg_price'] for h in portfolio_data)
            pl = total_value - invested
            pl_pct = (pl / invested * 100) if invested else 0
            
            slide2 = prs.slides.add_slide(prs.slide_layouts[6])
            bg2 = slide2.background.fill
            bg2.solid()
            bg2.fore_color.rgb = RGBColor(247, 249, 252)
            
            # Section title
            tb = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
            tfp = tb.text_frame.paragraphs[0]
            tfp.text = "KPI Portafoglio"
            tfp.font.size = Pt(28)
            tfp.font.bold = True
            tfp.font.color.rgb = RGBColor(10, 37, 64)
            
            kpi_data = [
                ('Valore Totale', f'EUR {total_value:,.2f}', RGBColor(10, 37, 64)),
                ('Investito', f'EUR {invested:,.2f}', RGBColor(138, 148, 166)),
                ('P&L', f'EUR {pl:,.2f} ({pl_pct:.1f}%)', RGBColor(15, 123, 63) if pl >= 0 else RGBColor(180, 35, 24)),
            ]
            
            for i, (label, value, color) in enumerate(kpi_data):
                x = 0.5 + i * 4.2
                box = slide2.shapes.add_shape(1, Inches(x), Inches(1.2), Inches(3.8), Inches(1.8))
                box.fill.solid()
                box.fill.fore_color.rgb = RGBColor(255, 255, 255)
                box.line.color.rgb = RGBColor(232, 236, 240)
                
                lbl = slide2.shapes.add_textbox(Inches(x + 0.2), Inches(1.4), Inches(3.4), Inches(0.4))
                lp = lbl.text_frame.paragraphs[0]
                lp.text = label
                lp.font.size = Pt(14)
                lp.font.color.rgb = RGBColor(138, 148, 166)
                
                val = slide2.shapes.add_textbox(Inches(x + 0.2), Inches(1.9), Inches(3.4), Inches(0.6))
                vp = val.text_frame.paragraphs[0]
                vp.text = value
                vp.font.size = Pt(24)
                vp.font.bold = True
                vp.font.color.rgb = color
            
            # Holdings table slide
            slide3 = prs.slides.add_slide(prs.slide_layouts[6])
            bg3 = slide3.background.fill
            bg3.solid()
            bg3.fore_color.rgb = RGBColor(247, 249, 252)
            
            ht = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
            htp = ht.text_frame.paragraphs[0]
            htp.text = "Posizioni"
            htp.font.size = Pt(28)
            htp.font.bold = True
            htp.font.color.rgb = RGBColor(10, 37, 64)
            
            rows = len(portfolio_data) + 1
            cols = 5
            table = slide3.shapes.add_table(rows, cols, Inches(0.5), Inches(1.1), Inches(12), Inches(0.6 * rows)).table
            
            headers = ['ETF', 'Quote', 'Prezzo Medio', 'Valore', 'P&L']
            for i, h in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = h
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(10, 37, 64)
            
            for r, h in enumerate(portfolio_data):
                cp = h.get('current_price', h['avg_price'])
                val = h['shares'] * cp
                plh = h['shares'] * (cp - h['avg_price'])
                row_data = [h['ticker'], f"{h['shares']:.1f}", f"{h['avg_price']:.2f}", f"{val:,.2f}", f"{plh:,.2f}"]
                for c, text in enumerate(row_data):
                    cell = table.cell(r + 1, c)
                    cell.text = text
                    cell.text_frame.paragraphs[0].font.size = Pt(11)
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(45, 58, 74)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if r % 2 == 0 else RGBColor(247, 249, 252)
        
        # Save
        pptx_path = os.path.join(UPLOAD_FOLDER, f'report_{user_id}_{int(datetime.now().timestamp())}.pptx')
        prs.save(pptx_path)
        return pptx_path
    except ImportError:
        return None

def build_report_html(user_id, portfolio_data, report_type):
    user = get_current_user() or {'name': 'Investitore'}
    total_value = sum(h['shares'] * h.get('current_price', h['avg_price']) for h in portfolio_data) if portfolio_data else 0
    invested = sum(h['shares'] * h['avg_price'] for h in portfolio_data) if portfolio_data else 0
    pl = total_value - invested
    
    rows = ""
    for h in portfolio_data:
        cp = h.get('current_price', h['avg_price'])
        val = h['shares'] * cp
        plh = h['shares'] * (cp - h['avg_price'])
        pl_color = "#0F7B3F" if plh >= 0 else "#B42318"
        rows += f"""
        <tr>
            <td style="padding:10px 12px;border-bottom:1px solid #E8ECF0;font-weight:600;">{h['ticker']}</td>
            <td style="padding:10px 12px;border-bottom:1px solid #E8ECF0;text-align:right;">{h['shares']:.2f}</td>
            <td style="padding:10px 12px;border-bottom:1px solid #E8ECF0;text-align:right;">{h['avg_price']:.2f} EUR</td>
            <td style="padding:10px 12px;border-bottom:1px solid #E8ECF0;text-align:right;">{cp:.2f} EUR</td>
            <td style="padding:10px 12px;border-bottom:1px solid #E8ECF0;text-align:right;font-weight:600;">{val:,.2f} EUR</td>
            <td style="padding:10px 12px;border-bottom:1px solid #E8ECF0;text-align:right;color:{pl_color};font-weight:600;">{plh:,.2f} EUR</td>
        </tr>"""
    
    return f"""
    <!DOCTYPE html>
    <html>
    <head>
        <meta charset="UTF-8">
        <style>
            body {{ font-family: 'Inter', -apple-system, sans-serif; margin: 0; padding: 40px; background: #fff; color: #0A2540; }}
            .header {{ background: #0A2540; color: #fff; padding: 30px 40px; margin: -40px -40px 30px; }}
            .header h1 {{ margin: 0; font-size: 28px; letter-spacing: -0.02em; }}
            .header p {{ margin: 8px 0 0; opacity: 0.7; font-size: 14px; }}
            .kpi-grid {{ display: grid; grid-template-columns: repeat(3, 1fr); gap: 16px; margin-bottom: 30px; }}
            .kpi-card {{ background: #F7F9FC; border: 1px solid #E8ECF0; border-radius: 12px; padding: 20px; }}
            .kpi-label {{ font-size: 12px; color: #8A94A6; font-weight: 600; text-transform: uppercase; margin-bottom: 6px; }}
            .kpi-value {{ font-size: 22px; font-weight: 700; color: #0A2540; }}
            table {{ width: 100%; border-collapse: collapse; font-size: 13px; }}
            th {{ text-align: left; padding: 12px; background: #0A2540; color: #fff; font-size: 11px; font-weight: 600; text-transform: uppercase; }}
            td {{ padding: 10px 12px; }}
            .footer {{ margin-top: 40px; padding-top: 20px; border-top: 1px solid #E8ECF0; font-size: 11px; color: #8A94A6; text-align: center; }}
        </style>
    </head>
    <body>
        <div class="header">
            <h1>Riepilogo Portafoglio ETF</h1>
            <p>{user['name']} &middot; {datetime.now().strftime('%d/%m/%Y %H:%M')}</p>
        </div>
        
        <div class="kpi-grid">
            <div class="kpi-card">
                <div class="kpi-label">Valore Totale</div>
                <div class="kpi-value">EUR {total_value:,.2f}</div>
            </div>
            <div class="kpi-card">
                <div class="kpi-label">Investito</div>
                <div class="kpi-value">EUR {invested:,.2f}</div>
            </div>
            <div class="kpi-card">
                <div class="kpi-label">P&L</div>
                <div class="kpi-value" style="color:{'#0F7B3F' if pl >= 0 else '#B42318'}">EUR {pl:,.2f}</div>
            </div>
        </div>
        
        <h2 style="font-size:18px;color:#0A2540;margin-bottom:16px;">Posizioni</h2>
        <table>
            <thead>
                <tr>
                    <th>ETF</th>
                    <th style="text-align:right;">Quote</th>
                    <th style="text-align:right;">Prezzo Medio</th>
                    <th style="text-align:right;">Prezzo Attuale</th>
                    <th style="text-align:right;">Valore</th>
                    <th style="text-align:right;">P&L</th>
                </tr>
            </thead>
            <tbody>{rows}</tbody>
        </table>
        
        <div class="footer">
            Generato da ETF Tracker &middot; Questo report non costituisce consulenza finanziaria
        </div>
    </body>
    </html>
    """

# ═══════════════════════════════════════════
# ROUTES
# ═══════════════════════════════════════════

@app.route('/')
def index():
    return render_template('index.html')

# ─── AUTH ───
@app.route('/api/auth/register', methods=['POST'])
def register():
    data = request.get_json() or {}
    email = data.get('email', '').strip().lower()
    password = data.get('password', '')
    name = data.get('name', 'Investitore').strip()
    
    if not email or not password or len(password) < 8:
        return jsonify({'error': 'Email e password (min 8 caratteri) richiesti'}), 400
    
    db = get_db()
    cursor = db.cursor()
    
    cursor.execute('SELECT id FROM users WHERE email = ?', (email,))
    if cursor.fetchone():
        return jsonify({'error': 'Email già registrata'}), 409
    
    initials = ''.join(w[0] for w in name.split() if w).upper()[:2] or 'IN'
    cursor.execute('''
        INSERT INTO users (email, password_hash, name, initials)
        VALUES (?, ?, ?, ?)
    ''', (email, hash_password(password), name, initials))
    db.commit()
    
    user_id = cursor.lastrowid
    token = generate_token(user_id)
    
    return jsonify({'token': token, 'user': {'id': user_id, 'email': email, 'name': name, 'initials': initials}}), 201

@app.route('/api/auth/login', methods=['POST'])
def login():
    data = request.get_json() or {}
    email = data.get('email', '').strip().lower()
    password = data.get('password', '')
    
    if not email or not password:
        return jsonify({'error': 'Email e password richiesti'}), 400
    
    db = get_db()
    cursor = db.cursor()
    cursor.execute('SELECT * FROM users WHERE email = ?', (email,))
    user = cursor.fetchone()
    
    if not user or user['password_hash'] != hash_password(password):
        return jsonify({'error': 'Credenziali non valide'}), 401
    
    cursor.execute('UPDATE users SET last_login = CURRENT_TIMESTAMP WHERE id = ?', (user['id'],))
    db.commit()
    
    token = generate_token(user['id'])
    return jsonify({
        'token': token,
        'user': {'id': user['id'], 'email': user['email'], 'name': user['name'], 'initials': user['initials']}
    })

@app.route('/api/auth/me', methods=['GET'])
@token_required
def get_me():
    user = get_current_user()
    if not user:
        return jsonify({'error': 'Utente non trovato'}), 404
    return jsonify({'id': user['id'], 'email': user['email'], 'name': user['name'], 'initials': user['initials']})

@app.route('/api/auth/logout', methods=['POST'])
@token_required
def logout():
    return jsonify({'success': True})

@app.route('/api/auth/onboarding', methods=['POST'])
@token_required
def save_onboarding():
    data = request.get_json() or {}
    db = get_db()
    cursor = db.cursor()
    cursor.execute('''
        UPDATE users SET onboarding_done = 1, onboarding_answers = ? WHERE id = ?
    ''', (json.dumps(data.get('answers', {})), g.current_user_id))
    db.commit()
    return jsonify({'success': True})

# ─── PORTFOLIOS ───
@app.route('/api/portfolios', methods=['GET', 'POST'])
@token_required
def portfolios():
    db = get_db()
    cursor = db.cursor()
    
    if request.method == 'GET':
        cursor.execute('SELECT * FROM portfolios WHERE user_id = ?', (g.current_user_id,))
        portfolios = [dict(row) for row in cursor.fetchall()]
        for p in portfolios:
            cursor.execute('SELECT * FROM holdings WHERE portfolio_id = ?', (p['id'],))
            p['holdings'] = [dict(row) for row in cursor.fetchall()]
        return jsonify(portfolios)
    
    data = request.get_json() or {}
    name = data.get('name', 'Portafoglio Principale')
    strategy = data.get('strategy', 'custom')
    
    cursor.execute('INSERT INTO portfolios (user_id, name, strategy) VALUES (?, ?, ?)',
                   (g.current_user_id, name, strategy))
    db.commit()
    return jsonify({'id': cursor.lastrowid, 'name': name, 'strategy': strategy}), 201

@app.route('/api/portfolios/<int:pid>/holdings', methods=['POST'])
@token_required
def add_holding(pid):
    data = request.get_json() or {}
    db = get_db()
    cursor = db.cursor()
    
    cursor.execute('SELECT * FROM portfolios WHERE id = ? AND user_id = ?', (pid, g.current_user_id))
    if not cursor.fetchone():
        return jsonify({'error': 'Portafoglio non trovato'}), 404
    
    cursor.execute('''
        INSERT INTO holdings (portfolio_id, isin, ticker, name, shares, avg_price)
        VALUES (?, ?, ?, ?, ?, ?)
    ''', (pid, data.get('isin'), data.get('ticker'), data.get('name'),
          data.get('shares', 0), data.get('avgPrice', 0)))
    db.commit()
    return jsonify({'id': cursor.lastrowid}), 201

@app.route('/api/holdings/<int:hid>', methods=['DELETE'])
@token_required
def delete_holding(hid):
    db = get_db()
    cursor = db.cursor()
    cursor.execute('''
        DELETE FROM holdings WHERE id = ? AND portfolio_id IN 
        (SELECT id FROM portfolios WHERE user_id = ?)
    ''', (hid, g.current_user_id))
    db.commit()
    return jsonify({'success': True})

# ─── ALERTS ───
@app.route('/api/alerts', methods=['GET', 'POST'])
@token_required
def alerts():
    db = get_db()
    cursor = db.cursor()
    
    if request.method == 'GET':
        cursor.execute('SELECT * FROM alerts WHERE user_id = ?', (g.current_user_id,))
        return jsonify([dict(row) for row in cursor.fetchall()])
    
    data = request.get_json() or {}
    cursor.execute('''
        INSERT INTO alerts (user_id, etf_isin, etf_ticker, alert_type, threshold, current_price, channels)
        VALUES (?, ?, ?, ?, ?, ?, ?)
    ''', (g.current_user_id, data.get('etfIsin'), data.get('etfTicker'),
          data.get('type', 'below'), data.get('threshold', 0),
          data.get('currentPrice', 0), json.dumps(data.get('channels', ['email', 'push']))))
    db.commit()
    return jsonify({'id': cursor.lastrowid}), 201

@app.route('/api/alerts/<int:aid>', methods=['PATCH', 'DELETE'])
@token_required
def alert_detail(aid):
    db = get_db()
    cursor = db.cursor()
    
    if request.method == 'DELETE':
        cursor.execute('DELETE FROM alerts WHERE id = ? AND user_id = ?', (aid, g.current_user_id))
        db.commit()
        return jsonify({'success': True})
    
    data = request.get_json() or {}
    cursor.execute('UPDATE alerts SET active = ? WHERE id = ? AND user_id = ?',
                   (1 if data.get('active') else 0, aid, g.current_user_id))
    db.commit()
    return jsonify({'success': True})

# ─── SETTINGS ───
@app.route('/api/settings', methods=['GET', 'PUT'])
@token_required
def settings():
    db = get_db()
    cursor = db.cursor()
    
    if request.method == 'GET':
        cursor.execute('SELECT key, value FROM settings WHERE user_id = ?', (g.current_user_id,))
        return jsonify({row['key']: row['value'] for row in cursor.fetchall()})
    
    data = request.get_json() or {}
    for key, value in data.items():
        cursor.execute('''
            INSERT INTO settings (user_id, key, value) VALUES (?, ?, ?)
            ON CONFLICT(user_id, key) DO UPDATE SET value = excluded.value
        ''', (g.current_user_id, key, json.dumps(value) if isinstance(value, (dict, list)) else str(value)))
    db.commit()
    return jsonify({'success': True})

# ─── REPORTS ───
@app.route('/api/reports/generate', methods=['POST'])
@token_required
def generate_report():
    data = request.get_json() or {}
    report_type = data.get('type', 'portfolio_summary')
    report_format = data.get('format', 'pdf')
    
    # Get user's portfolio data
    db = get_db()
    cursor = db.cursor()
    cursor.execute('SELECT * FROM portfolios WHERE user_id = ?', (g.current_user_id,))
    portfolios = [dict(row) for row in cursor.fetchall()]
    
    holdings = []
    for p in portfolios:
        cursor.execute('SELECT * FROM holdings WHERE portfolio_id = ?', (p['id'],))
        for h in cursor.fetchall():
            hdict = dict(h)
            # Add current price from static data
            from_static = get_etf_data(hdict['isin'])
            hdict['current_price'] = from_static.get('price', hdict['avg_price'])
            hdict['chg1d'] = from_static.get('chg1d', 0)
            hdict['chg1y'] = from_static.get('chg1y', 0)
            holdings.append(hdict)
    
    file_path = None
    if report_format == 'pdf':
        file_path = generate_pdf_report(g.current_user_id, holdings, report_type)
    elif report_format == 'pptx':
        file_path = generate_pptx_report(g.current_user_id, holdings, report_type)
    
    if not file_path:
        return jsonify({'error': 'Generazione report fallita'}), 500
    
    cursor.execute('''
        INSERT INTO reports (user_id, report_type, format, title, content_json, file_path)
        VALUES (?, ?, ?, ?, ?, ?)
    ''', (g.current_user_id, report_type, report_format,
          f'Riepilogo {report_type}', json.dumps({'holdings_count': len(holdings)}), file_path))
    db.commit()
    
    return jsonify({
        'id': cursor.lastrowid,
        'file_path': file_path,
        'download_url': f'/api/reports/download/{cursor.lastrowid}'
    })

@app.route('/api/reports/download/<int:rid>')
@token_required
def download_report(rid):
    db = get_db()
    cursor = db.cursor()
    cursor.execute('SELECT * FROM reports WHERE id = ? AND user_id = ?', (rid, g.current_user_id))
    report = cursor.fetchone()
    if not report:
        return jsonify({'error': 'Report non trovato'}), 404
    
    report = dict(report)
    if not os.path.exists(report['file_path']):
        return jsonify({'error': 'File non trovato'}), 404
    
    ext = os.path.splitext(report['file_path'])[1]
    mime = 'application/pdf' if ext == '.pdf' else 'application/vnd.openxmlformats-officedocument.presentationml.presentation' if ext == '.pptx' else 'application/octet-stream'
    return send_file(report['file_path'], mimetype=mime, as_download=True,
                     download_name=f"ETF_Report_{report['report_type']}{ext}")

@app.route('/api/reports/email', methods=['POST'])
@token_required
def email_report():
    data = request.get_json() or {}
    report_format = data.get('format', 'pdf')
    
    # Generate report
    db = get_db()
    cursor = db.cursor()
    cursor.execute('SELECT * FROM portfolios WHERE user_id = ?', (g.current_user_id,))
    portfolios = [dict(row) for row in cursor.fetchall()]
    
    holdings = []
    for p in portfolios:
        cursor.execute('SELECT * FROM holdings WHERE portfolio_id = ?', (p['id'],))
        for h in cursor.fetchall():
            hdict = dict(h)
            from_static = get_etf_data(hdict['isin'])
            hdict['current_price'] = from_static.get('price', hdict['avg_price'])
            holdings.append(hdict)
    
    user = get_current_user()
    if not user:
        return jsonify({'error': 'Utente non trovato'}), 404
    
    # Generate file
    file_path = None
    if report_format == 'pdf':
        file_path = generate_pdf_report(g.current_user_id, holdings, 'portfolio_summary')
    elif report_format == 'pptx':
        file_path = generate_pptx_report(g.current_user_id, holdings, 'portfolio_summary')
    
    # Save report record
    cursor.execute('''
        INSERT INTO reports (user_id, report_type, format, title, content_json, file_path)
        VALUES (?, ?, ?, ?, ?, ?)
    ''', (g.current_user_id, 'portfolio_summary', report_format,
          'Riepilogo Portafoglio', json.dumps({'holdings_count': len(holdings)}), file_path))
    db.commit()
    report_id = cursor.lastrowid
    
    # Build email
    total_value = sum(h['shares'] * h.get('current_price', h['avg_price']) for h in holdings)
    invested = sum(h['shares'] * h['avg_price'] for h in holdings)
    pl = total_value - invested
    
    email_html = f"""
    <div style="font-family:Inter,sans-serif;max-width:600px;margin:0 auto;padding:24px;">
        <div style="background:#0A2540;color:#fff;padding:24px;border-radius:12px 12px 0 0;">
            <h2 style="margin:0;font-size:22px;">ETF Tracker - Riepilogo Portafoglio</h2>
            <p style="margin:8px 0 0;opacity:0.7;font-size:14px;">{user['name']} &middot; {datetime.now().strftime('%d/%m/%Y')}</p>
        </div>
        <div style="background:#fff;border:1px solid #E8ECF0;border-top:none;padding:24px;border-radius:0 0 12px 12px;">
            <div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:12px;margin-bottom:24px;">
                <div style="background:#F7F9FC;padding:16px;border-radius:8px;text-align:center;">
                    <p style="font-size:11px;color:#8A94A6;margin:0 0 4px;">VALORE TOTALE</p>
                    <p style="font-size:18px;font-weight:700;color:#0A2540;margin:0;">EUR {total_value:,.2f}</p>
                </div>
                <div style="background:#F7F9FC;padding:16px;border-radius:8px;text-align:center;">
                    <p style="font-size:11px;color:#8A94A6;margin:0 0 4px;">INVESTITO</p>
                    <p style="font-size:18px;font-weight:700;color:#0A2540;margin:0;">EUR {invested:,.2f}</p>
                </div>
                <div style="background:#F7F9FC;padding:16px;border-radius:8px;text-align:center;">
                    <p style="font-size:11px;color:#8A94A6;margin:0 0 4px;">P&L</p>
                    <p style="font-size:18px;font-weight:700;margin:0;color:{'#0F7B3F' if pl >= 0 else '#B42318'};">EUR {pl:,.2f}</p>
                </div>
            </div>
            <p style="font-size:13px;color:#8A94A6;text-align:center;">
                Il report completo e allegato. Accedi a <a href="#" style="color:#0F7B3F;">ETF Tracker</a> per analisi dettagliate.
            </p>
        </div>
    </div>
    """
    
    success = send_email(user['email'], 'ETF Tracker - Riepilogo Portafoglio', email_html,
                        attachments=[file_path] if file_path else None)
    
    if success:
        cursor.execute('UPDATE reports SET status = ? WHERE id = ?', ('sent', report_id))
        db.commit()
        return jsonify({'success': True, 'message': 'Email inviata con successo'})
    else:
        return jsonify({'success': False, 'message': 'Email non inviata - controlla la configurazione SMTP'}), 200

# ─── COMMUNITY ───
@app.route('/api/community/posts', methods=['GET', 'POST'])
@token_required
def community_posts():
    db = get_db()
    cursor = db.cursor()
    
    if request.method == 'GET':
        cursor.execute('SELECT * FROM community_posts ORDER BY created_at DESC')
        return jsonify([dict(row) for row in cursor.fetchall()])
    
    data = request.get_json() or {}
    user = get_current_user()
    cursor.execute('''
        INSERT INTO community_posts (user_id, user_name, avatar, content, portfolio_snapshot)
        VALUES (?, ?, ?, ?, ?)
    ''', (g.current_user_id, user['name'], user['initials'],
          data.get('content', ''), json.dumps(data.get('portfolio'))))
    db.commit()
    return jsonify({'id': cursor.lastrowid}), 201

@app.route('/api/community/posts/<int:pid>/like', methods=['POST'])
@token_required
def like_post(pid):
    db = get_db()
    cursor = db.cursor()
    cursor.execute('UPDATE community_posts SET likes = likes + 1 WHERE id = ?', (pid,))
    db.commit()
    return jsonify({'success': True})

# ─── STATIC DATA ───
@app.route('/api/data/etfs')
def get_etfs():
    return jsonify(ETF_DATABASE)

@app.route('/api/data/models')
def get_models():
    return jsonify(MODEL_PORTFOLIOS)

def get_etf_data(isin):
    for e in ETF_DATABASE:
        if e['isin'] == isin:
            return e
    return {}

# ═══════════════════════════════════════════
# STATIC DATA (inline per plug-and-play)
# ═══════════════════════════════════════════

ETF_DATABASE = [
    {"isin":"IE00B4L5Y983","ticker":"VWCE","name":"Vanguard FTSE All-World","price":118.50,"chg1d":0.45,"chg1y":12.3,"ter":0.22,"aum":1850000,"asset":"Azionario","region":"Globale","replication":"Fisica","distribution":"Accumulazione","domicile":"Irlanda","issuer":"Vanguard","rating":5,"chg5y":52.0},
    {"isin":"IE00B4L5YY86","ticker":"IWDA","name":"iShares Core MSCI World","price":85.20,"chg1d":0.32,"chg1y":14.1,"ter":0.20,"aum":980000,"asset":"Azionario","region":"Globale","replication":"Fisica","distribution":"Accumulazione","domicile":"Irlanda","issuer":"BlackRock","rating":5,"chg5y":55.0},
    {"isin":"IE00BDBRDM35","ticker":"EUNL","name":"iShares Core Euro STOXX 50","price":42.80,"chg1d":0.18,"chg1y":8.5,"ter":0.10,"aum":125000,"asset":"Azionario","region":"Europa","replication":"Fisica","distribution":"Accumulazione","domicile":"Irlanda","issuer":"BlackRock","rating":4,"chg5y":28.0},
    {"isin":"IE00BDBRDM43","ticker":"EXS1","name":"iShares STOXX Europe 600","price":95.40,"chg1d":0.25,"chg1y":9.2,"ter":0.20,"aum":85000,"asset":"Azionario","region":"Europa","replication":"Fisica","distribution":"Accumulazione","domicile":"Irlanda","issuer":"BlackRock","rating":4,"chg5y":30.0},
    {"isin":"IE00BKM4GZ66","ticker":"EMAE","name":"iShares Core MSCI EM IMI","price":32.15,"chg1d":0.85,"chg1y":6.8,"ter":0.18,"aum":145000,"asset":"Azionario","region":"Emergenti","replication":"Fisica","distribution":"Accumulazione","domicile":"Irlanda","issuer":"BlackRock","rating":4,"chg5y":18.0},
    {"isin":"IE00BDBRDM51","ticker":"IUSN","name":"iShares MSCI World Small Cap","price":8.95,"chg1d":0.42,"chg1y":11.5,"ter":0.35,"aum":42000,"asset":"Azionario","region":"Globale","replication":"Fisica","distribution":"Accumulazione","domicile":"Irlanda","issuer":"BlackRock","rating":4,"chg5y":35.0},
    {"isin":"IE00BDBRDM69","ticker":"IBGL","name":"iShares Euro Government Bond","price":136.20,"chg1d":-0.15,"chg1y":-2.1,"ter":0.09,"aum":185000,"asset":"Obbligazionario","region":"Europa","replication":"Fisica","distribution":"Distribuzione","domicile":"Irlanda","issuer":"BlackRock","rating":4,"chg5y":-8.0},
    {"isin":"IE00BDBRDM77","ticker":"IEGA","name":"iShares Core Euro Gov Bond","price":128.50,"chg1d":-0.12,"chg1y":-1.8,"ter":0.12,"aum":95000,"asset":"Obbligazionario","region":"Europa","replication":"Fisica","distribution":"Accumulazione","domicile":"Irlanda","issuer":"BlackRock","rating":4,"chg5y":-6.0},
    {"isin":"IE00BDBRDM85","ticker":"IBCI","name":"iShares Inflation Linked Govt","price":68.40,"chg1d":0.08,"chg1y":1.2,"ter":0.10,"aum":32000,"asset":"Obbligazionario","region":"Europa","replication":"Fisica","distribution":"Distribuzione","domicile":"Irlanda","issuer":"BlackRock","rating":4,"chg5y":5.0},
    {"isin":"IE00BDBRDM93","ticker":"XAD1","name":"Xtrackers MSCI AC World","price":98.75,"chg1d":0.38,"chg1y":13.5,"ter":0.19,"aum":65000,"asset":"Azionario","region":"Globale","replication":"Fisica","distribution":"Accumulazione","domicile":"Irlanda","issuer":"DWS","rating":4,"chg5y":50.0},
    {"isin":"IE00BDBRDM01","ticker":"XBLC","name":"Xtrackers MSCI World Momentum","price":62.30,"chg1d":0.55,"chg1y":16.2,"ter":0.30,"aum":28000,"asset":"Azionario","region":"Globale","replication":"Campionamento","distribution":"Accumulazione","domicile":"Irlanda","issuer":"DWS","rating":4,"chg5y":58.0},
    {"isin":"IE00BDBRDM19","ticker":"Xetra-Gold","name":"Xetra-Gold","price":48.90,"chg1d":0.22,"chg1y":18.5,"ter":0.36,"aum":7800,"asset":"Commodities","region":"Globale","replication":"Fisica","distribution":"Accumulazione","domicile":"Germania","issuer":"Deutsche Borse","rating":3,"chg5y":42.0},
    {"isin":"IE00BDBRDM27","ticker":"IUSQ","name":"iShares MSCI ACWI","price":92.40,"chg1d":0.40,"chg1y":12.8,"ter":0.20,"aum":89000,"asset":"Azionario","region":"Globale","replication":"Fisica","distribution":"Accumulazione","domicile":"Irlanda","issuer":"BlackRock","rating":5,"chg5y":51.0},
]

MODEL_PORTFOLIOS = [
    {"id":"all-weather","name":"All-Weather","author":"Ray Dalio","risk":"Medio","riskLevel":3,"cagr":7.8,"maxDD":"-14.5","sharpe":0.72,"philosophy":"Asset allocation che performa in ogni condizione macroeconomica","allocation":[{"name":"Azionario","value":30,"color":"#0A2540"},{"name":"Obbligazionario","value":40,"color":"#1E5AA0"},{"name":"Oro","value":15,"color":"#D69E2E"},{"name":"Commodities","value":15,"color":"#68A063"}]},
    {"id":"bogleheads","name":"Bogleheads 3-Fund","author":"John Bogle","risk":"Medio","riskLevel":3,"cagr":8.2,"maxDD":"-18.3","sharpe":0.68,"philosophy":"Semplicità totale: azionario globale + obbligazionario + emergenti","allocation":[{"name":"Azionario Global","value":60,"color":"#0A2540"},{"name":"Emergenti","value":20,"color":"#1E5AA0"},{"name":"Obbligazionario","value":20,"color":"#8A94A6"}]},
    {"id":"permanent","name":"Permanent Portfolio","author":"Harry Browne","risk":"Basso","riskLevel":2,"cagr":5.4,"maxDD":"-8.1","sharpe":0.55,"philosophy":"Equa ripartizione per proteggere da ogni scenario economico","allocation":[{"name":"Azionario","value":25,"color":"#0A2540"},{"name":"Obbligazionario","value":25,"color":"#1E5AA0"},{"name":"Oro","value":25,"color":"#D69E2E"},{"name":"Cash","value":25,"color":"#8A94A6"}]},
    {"id":"coffeehouse","name":"Coffeehouse Portfolio","author":"Bill Schultheis","risk":"Basso","riskLevel":2,"cagr":6.8,"maxDD":"-12.4","sharpe":0.62,"philosophy":"Diversificazione estrema con 10 asset class","allocation":[{"name":"Large Cap","value":10,"color":"#0A2540"},{"name":"Large Value","value":10,"color":"#1E3A5F"},{"name":"Small Cap","value":10,"color":"#1E5AA0"},{"name":"Small Value","value":10,"color":"#0F7B3F"},{"name":"REITs","value":10,"color":"#68A063"},{"name":"Gov Bond","value":10,"color":"#D69E2E"},{"name":"Corp Bond","value":10,"color":"#8A94A6"},{"name":"TIPS","value":10,"color":"#4A5568"},{"name":"Emergenti","value":10,"color":"#2D3748"},{"name":"EAFE","value":10,"color":"#061528"}]},
    {"id":"golden","name":"Golden Butterfly","author":"Tyler","risk":"Medio-Basso","riskLevel":2,"cagr":7.1,"maxDD":"-11.2","sharpe":0.70,"philosophy":"Variante Permanent con tilt value e small cap","allocation":[{"name":"S&P 500","value":20,"color":"#0A2540"},{"name":"Small Cap Value","value":20,"color":"#1E3A5F"},{"name":"Long Gov Bond","value":20,"color":"#1E5AA0"},{"name":"Short Gov Bond","value":20,"color":"#8A94A6"},{"name":"Oro","value":20,"color":"#D69E2E"}]},
    {"id":"growth","name":"100% Growth","author":"ETF Tracker","risk":"Alto","riskLevel":4,"cagr":10.5,"maxDD":"-28.7","sharpe":0.58,"philosophy":"Massima esposizione azionaria per obiettivi di crescita","allocation":[{"name":"Azionario Global","value":70,"color":"#0A2540"},{"name":"Small Cap","value":20,"color":"#1E3A5F"},{"name":"Emergenti","value":10,"color":"#0F7B3F"}]},
]

CORRELATION_MATRIX = [
    {"etf":"VWCE","VWCE":1.0,"IWDA":0.96,"EUNL":0.88,"EXS1":0.90,"EMAE":0.82,"IUSN":0.93,"IBGL":0.15,"IEGA":0.18,"IBCI":0.22,"XAD1":0.97,"XBLC":0.95,"Xetra-Gold":-0.08,"IUSQ":0.98},
    {"etf":"IWDA","VWCE":0.96,"IWDA":1.0,"EUNL":0.85,"EXS1":0.87,"EMAE":0.78,"IUSN":0.91,"IBGL":0.12,"IEGA":0.15,"IBCI":0.18,"XAD1":0.99,"XBLC":0.96,"Xetra-Gold":-0.12,"IUSQ":0.97},
    {"etf":"EUNL","VWCE":0.88,"IWDA":0.85,"EUNL":1.0,"EXS1":0.95,"EMAE":0.65,"IUSN":0.78,"IBGL":0.25,"IEGA":0.28,"IBCI":0.30,"XAD1":0.86,"XBLC":0.82,"Xetra-Gold":-0.05,"IUSQ":0.88},
    {"etf":"EXS1","VWCE":0.90,"IWDA":0.87,"EUNL":0.95,"EXS1":1.0,"EMAE":0.70,"IUSN":0.80,"IBGL":0.22,"IEGA":0.25,"IBCI":0.28,"XAD1":0.88,"XBLC":0.84,"Xetra-Gold":-0.02,"IUSQ":0.90},
    {"etf":"EMAE","VWCE":0.82,"IWDA":0.78,"EUNL":0.65,"EXS1":0.70,"EMAE":1.0,"IUSN":0.75,"IBGL":-0.05,"IEGA":-0.02,"IBCI":0.05,"XAD1":0.80,"XBLC":0.78,"Xetra-Gold":0.10,"IUSQ":0.83},
    {"etf":"IUSN","VWCE":0.93,"IWDA":0.91,"EUNL":0.78,"EXS1":0.80,"EMAE":0.75,"IUSN":1.0,"IBGL":0.08,"IEGA":0.10,"IBCI":0.12,"XAD1":0.92,"XBLC":0.94,"Xetra-Gold":-0.10,"IUSQ":0.93},
    {"etf":"IBGL","VWCE":0.15,"IWDA":0.12,"EUNL":0.25,"EXS1":0.22,"EMAE":-0.05,"IUSN":0.08,"IBGL":1.0,"IEGA":0.95,"IBCI":0.88,"XAD1":0.15,"XBLC":0.10,"Xetra-Gold":0.35,"IUSQ":0.14},
    {"etf":"IEGA","VWCE":0.18,"IWDA":0.15,"EUNL":0.28,"EXS1":0.25,"EMAE":-0.02,"IUSN":0.10,"IBGL":0.95,"IEGA":1.0,"IBCI":0.85,"XAD1":0.18,"XBLC":0.14,"Xetra-Gold":0.32,"IUSQ":0.17},
    {"etf":"IBCI","VWCE":0.22,"IWDA":0.18,"EUNL":0.30,"EXS1":0.28,"EMAE":0.05,"IUSN":0.12,"IBGL":0.88,"IEGA":0.85,"IBCI":1.0,"XAD1":0.22,"XBLC":0.20,"Xetra-Gold":0.38,"IUSQ":0.23},
    {"etf":"XAD1","VWCE":0.97,"IWDA":0.99,"EUNL":0.86,"EXS1":0.88,"EMAE":0.80,"IUSN":0.92,"IBGL":0.15,"IEGA":0.18,"IBCI":0.22,"XAD1":1.0,"XBLC":0.97,"Xetra-Gold":-0.10,"IUSQ":0.99},
    {"etf":"XBLC","VWCE":0.95,"IWDA":0.96,"EUNL":0.82,"EXS1":0.84,"EMAE":0.78,"IUSN":0.94,"IBGL":0.10,"IEGA":0.14,"IBCI":0.20,"XAD1":0.97,"XBLC":1.0,"Xetra-Gold":-0.08,"IUSQ":0.96},
    {"etf":"Xetra-Gold","VWCE":-0.08,"IWDA":-0.12,"EUNL":-0.05,"EXS1":-0.02,"EMAE":0.10,"IUSN":-0.10,"IBGL":0.35,"IEGA":0.32,"IBCI":0.38,"XAD1":-0.10,"XBLC":-0.08,"Xetra-Gold":1.0,"IUSQ":-0.09},
    {"etf":"IUSQ","VWCE":0.98,"IWDA":0.97,"EUNL":0.88,"EXS1":0.90,"EMAE":0.83,"IUSN":0.93,"IBGL":0.14,"IEGA":0.17,"IBCI":0.23,"XAD1":0.99,"XBLC":0.96,"Xetra-Gold":-0.09,"IUSQ":1.0},
]

ACADEMY_COURSES = [
    {"id":"abc-etf","title":"L'ABC degli ETF","level":"Base","duration":"3h","lessons":8},
    {"id":"analisi-etf","title":"Analisi e Confronto","level":"Intermedio","duration":"4h","lessons":10},
    {"id":"strategie","title":"Strategie Portafoglio","level":"Avanzato","duration":"6h","lessons":12},
    {"id":"tassazione","title":"Tassazione e Dichiarazione","level":"Intermedio","duration":"2h","lessons":5},
    {"id":"rischio","title":"Risk Management","level":"Avanzato","duration":"4h","lessons":7},
    {"id":"etf-vs-fondi","title":"ETF vs Fondi vs Azioni","level":"Base","duration":"2h","lessons":4},
]

MANUAL_MODULES = [
    {"id":"mod0","title":"Modulo 0: Premessa","subtitle":"Prima di cominciare","icon":"book-open","color":"#E8EEF3","chapters":[
        {"id":"ch0-1","num":"0.1","title":"Per chi e questo manuale","readTime":"3 min"},
        {"id":"ch0-2","num":"0.2","title":"Come usare questa guida","readTime":"2 min"},
    ]},
    {"id":"mod1","title":"Modulo 1: Fondamenti","subtitle":"Dalle basi alla scelta","icon":"layers","color":"#E0EDFF","chapters":[
        {"id":"ch1-1","num":"1.1","title":"Cos'e un ETF","readTime":"5 min"},
        {"id":"ch1-2","num":"1.2","title":"Replica fisica vs sintetica","readTime":"4 min"},
        {"id":"ch1-3","num":"1.3","title":"TER, NAV e Tracking Error","readTime":"6 min"},
        {"id":"ch1-4","num":"1.4","title":"Scegliere il giusto ETF","readTime":"7 min"},
        {"id":"ch1-5","num":"1.5","title":"Dove acquistare","readTime":"5 min"},
    ]},
    {"id":"mod2","title":"Modulo 2: Pratica","subtitle":"Dalla teoria al portafoglio","icon":"pie-chart","color":"#E6F4EC","chapters":[
        {"id":"ch2-1","num":"2.1","title":"Costruire il primo portafoglio","readTime":"8 min"},
        {"id":"ch2-2","num":"2.2","title":"Asset Allocation","readTime":"10 min"},
        {"id":"ch2-3","num":"2.3","title":"PAC vs PIC","readTime":"6 min"},
        {"id":"ch2-4","num":"2.4","title":"Ribilanciamento","readTime":"7 min"},
        {"id":"ch2-5","num":"2.5","title":"Tassazione in Italia","readTime":"9 min"},
        {"id":"ch2-6","num":"2.6","title":"Dichiarazione redditi","readTime":"8 min"},
    ]},
]

CHAPTER_SUMMARIES = {
    "ch1-1": {"title":"Cos'e un ETF","keyPoints":["ETF = Exchange Traded Fund: fondo negoziato in borsa","Traccia un indice (es. MSCI World) in modo passivo","Diversificazione implicita in un singolo titolo","Costi bassi (TER 0.05%-0.50%) vs fondi attivi (1-2%)","Liquidita: compra/vendi in tempo reale come un azione"]},
    "ch1-2": {"title":"Replica fisica vs sintetica","keyPoints":["Replica fisica: detiene davvero i titoli dell indice","Replica sintetica: usa swap con una controparte","Campione: replica fisica su sottoinsieme di titoli","Fisica = trasparenza maggiore, costi leggermente superiori","Sintetica = tracking error minore ma rischio controparte"]},
    "ch1-3": {"title":"TER, NAV e Tracking Error","keyPoints":["TER = Total Expense Ratio: costo annuo sul patrimonio","NAV = Net Asset Value: valore rezzo del fondo","Tracking Error = scostamento tra rendimento ETF e indice","Premiums/Discounts: NAV vs prezzo di mercato","Costi impliciti: swap fee, lending revenue, tasse"]},
    "ch2-1": {"title":"Costruire il primo portafoglio","keyPoints":["Inizia con 1-3 ETF massimo per semplicita","Core-satellite: 70-80% core globale, 20-30% tematici","Ottieni esposizione a migliaia di titoli con pochi ETF","Considera la tua propensione al rischio e orizzonte","Investi solo capitale che non ti serve a breve termine"]},
    "ch2-2": {"title":"Asset Allocation","keyPoints":["Allocazione = come dividi il capitale tra asset class","Rule of 100: 100 - eta = % azionario (approssimativo)","Diversificazione tra geografie, settori, dimensioni","Rebalancing: riporta le % ai target a intervalli regolari","L'asset allocation spiega il 90% della performance"]},
}

# ─── Init on startup ───
with app.app_context():
    init_db()

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port)
